"""
텍스트 추출 모듈.
PDF, docx, pptx, txt, md, 코드 파일 등에서 텍스트 추출.
"""

from pathlib import Path
from typing import Optional

from app.core.config import get_settings
from app.core.logging_config import get_logger

logger = get_logger(__name__)

# 지원 확장자
SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".java",
    ".kt",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".cs",
    ".rb",
    ".php",
    ".swift",
    ".sql",
    ".html",
    ".css",
    ".scss",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".csv",
    ".xml",
    ".sh",
    ".bat",
    ".ps1",
}


class TextExtractor:
    """다양한 파일 형식에서 텍스트 추출."""

    def __init__(self) -> None:
        self.settings = get_settings()
        self._max_size = self.settings.max_file_size_mb * 1024 * 1024

    def can_extract(self, path: Path) -> bool:
        """해당 경로에서 추출 가능한지 여부."""
        if not path.is_file():
            return False
        try:
            if path.stat().st_size > self._max_size:
                return False
        except OSError:
            return False
        return path.suffix.lower() in SUPPORTED_EXTENSIONS or path.suffix.lower() in {
            ".pdf",
            ".docx",
            ".pptx",
        }

    def extract(self, path: Path, encoding: str = "utf-8") -> Optional[str]:
        """
        파일에서 텍스트 추출.
        Returns:
            추출된 텍스트 또는 None
        """
        path = Path(path)
        suffix = path.suffix.lower()

        try:
            if suffix in {".txt", ".md", ".markdown"} or suffix in SUPPORTED_EXTENSIONS:
                return self._extract_text_file(path, encoding)
            if suffix == ".pdf":
                return self._extract_pdf(path)
            if suffix == ".docx":
                return self._extract_docx(path)
            if suffix == ".pptx":
                return self._extract_pptx(path)
        except Exception as e:
            logger.warning("텍스트 추출 실패", path=str(path), error=str(e))
        return None

    def _extract_text_file(self, path: Path, encoding: str) -> Optional[str]:
        """일반 텍스트/코드 파일."""
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            try:
                return path.read_text(encoding="cp949")
            except Exception:
                return None

    def _extract_pdf(self, path: Path) -> Optional[str]:
        """PDF에서 텍스트 추출."""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(path)
            parts = []
            for page in reader.pages[:500]:  # 상위 500페이지만
                text = page.extract_text()
                if text:
                    parts.append(text)
            return "\n\n".join(parts) if parts else None
        except Exception as e:
            logger.warning("PDF 추출 실패", path=str(path), error=str(e))
            return None

    def _extract_docx(self, path: Path) -> Optional[str]:
        """DOCX에서 텍스트 추출."""
        try:
            from docx import Document

            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.warning("DOCX 추출 실패", path=str(path), error=str(e))
            return None

    def _extract_pptx(self, path: Path) -> Optional[str]:
        """PPTX에서 텍스트 추출."""
        try:
            from pptx import Presentation

            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts) if parts else None
        except Exception as e:
            logger.warning("PPTX 추출 실패", path=str(path), error=str(e))
            return None
